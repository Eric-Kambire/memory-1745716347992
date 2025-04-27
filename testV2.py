import requests
import os
import re
import time
import random
import math
from PyPDF2 import PdfReader
from pptx import Presentation
import textwrap
from fpdf import FPDF
from datetime import datetime, timedelta

# Configuration de l'API
API_KEY = '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'
API_URL = 'https://openrouter.ai/api/v1/chat/completions'
headers = {
    'Authorization': f'Bearer {API_KEY}',
    'Content-Type': 'application/json'
}

class SmartPDF(FPDF):
    """PDF intelligent avec mise en forme adaptative"""
    def header(self):
        self.set_font('Arial', 'B', 14)
        self.cell(0, 10, 'MemoryMesh - Résumé & Quiz', 0, 1, 'C')
        self.line(10, 20, 200, 20)
    
    def add_section(self, title, content):
        self.set_font('Arial', 'B', 12)
        self.cell(0, 10, title, 0, 1)
        self.set_font('Arial', '', 11)
        self.multi_cell(0, 7, content)
        self.ln(5)
    
    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')


def extract_text(file_path):
    """Extraction multi-format avec gestion d'erreur améliorée"""
    text = ""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            reader = PdfReader(file_path)
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            raise ValueError(f"Format {ext} non supporté")
        return text[:15000]  # Limite token
    except Exception as e:
        print(f"Erreur extraction: {str(e)}")
        return None


def analyze_content(text, action):
    """Fonction unifiée d'analyse avec DeepSeek"""
    prompts = {
        'summary': f"Résume ce document en 5 points essentiels en français:\n{text}",
        'quiz': f"""Génère un quiz basé sur ce contenu:
- 3 QCM (4 options)
- 2 Vrai/Faux
- 1 question ouverte
Format exigé:
[Type] Question...
Options: A) ... B) ... 
Réponse: ...
Explication: ...

Contenu:
{textwrap.shorten(text, width=10000, placeholder='...')}"""
    }
    
    data = {
        "model": "deepseek/deepseek-chat:free",
        "messages": [{"role": "user", "content": prompts[action]}],
        "temperature": 0.3 if action == 'summary' else 0.7
    }
    
    try:
        response = requests.post(API_URL, json=data, headers=headers, timeout=25)
        response.raise_for_status()
        return response.json()['choices'][0]['message']['content']
    except Exception as e:
        print(f"Erreur API ({action}): {str(e)}")
        return None


def parse_quiz(quiz_text):
    """Parse le texte du quiz en structure de données avec gestion de formats variés"""
    quiz_data = []
    if not quiz_text or not isinstance(quiz_text, str):
        return []
    questions_raw = re.split(r'\n\s*\n', quiz_text)
    for q_raw in questions_raw:
        if not q_raw.strip():
            continue
        question = {}
        # Détection du type
        if re.search(r'\[QCM\]|QCM', q_raw, re.IGNORECASE):
            question['type'] = 'qcm'
        elif re.search(r'\[Vrai/Faux\]|Vrai/Faux', q_raw, re.IGNORECASE):
            question['type'] = 'vf'
        elif re.search(r'\[Question ouverte\]', q_raw, re.IGNORECASE):
            question['type'] = 'open'
        else:
            question['type'] = 'open'
        try:
            # Question text
            question_pattern = r'(?:\[.*?\]|QCM|Vrai/Faux)[\s:.]*(.*?)(?:Options:|Réponse:|$)'
            m = re.search(question_pattern, q_raw, re.DOTALL | re.IGNORECASE)
            question['question'] = m.group(1).strip() if m else q_raw.strip().splitlines()[0]
            # Options
            if question['type'] == 'qcm':
                opts = {}
                options_match = re.search(r'Options:(.*?)(?:Réponse:|$)', q_raw, re.DOTALL)
                opts_text = options_match.group(1).strip() if options_match else q_raw
                for opt in re.finditer(r'([A-D])[).]([^A-D]+)', opts_text):
                    opts[opt.group(1)] = opt.group(2).strip()
                question['options'] = opts or {'A':'Option A','B':'Option B','C':'Option C','D':'Option D'}
            elif question['type'] == 'vf':
                question['options'] = {'A':'Vrai','B':'Faux'}
            # Réponse
            ans_m = re.search(r'(?:Réponse|Answer)[^A-Za-z]*([A-D]|Vrai|Faux)', q_raw, re.IGNORECASE)
            question['answer'] = ans_m.group(1).strip() if ans_m else ('A' if question['type']=='qcm' else 'Pas de réponse fournie')
            # Explication
            exp_m = re.search(r'(?:Explication|Explanation)[^A-Za-z]*(.*?)$', q_raw, re.DOTALL | re.IGNORECASE)
            question['explanation'] = exp_m.group(1).strip() if exp_m else ''
        except Exception as e:
            question = {
                'type': 'open',
                'question': f"Parsing error: {q_raw[:50]}...",
                'options': {},
                'answer': 'Erreur',
                'explanation': str(e)
            }
        quiz_data.append(question)
    return quiz_data


def present_quiz(quiz_data):
    """Présente le quiz et enregistre les résultats"""
    quiz_results, total_time = [], 0
    if not quiz_data:
        print("Aucune question.")
        return [], 0
    print("===== QUIZ INTERACTIF =====")
    for idx, q in enumerate(quiz_data, 1):
        print(f"\nQuestion {idx}: {q['question']}")
        if q['type']=='qcm':
            for k,v in q['options'].items(): print(f" {k}) {v}")
        elif q['type']=='vf': print(" A) Vrai\n B) Faux")
        start = time.time()
        ans = input("Votre réponse: ").strip().upper()
        duration = time.time() - start
        total_time += duration
        user_ans = ans[0] if ans and ans[0] in "ABCD" else ans
        correct = None
        if q['type'] in ['qcm','vf']:
            corr = re.search(r'([A-D])', str(q['answer']))
            correct = user_ans == (corr.group(1) if corr else 'A')
            corr_norm = corr.group(1) if corr else 'A'
        else:
            corr_norm = None
        quiz_results.append({
            'question': q['question'],
            'type': q['type'],
            'user_answer': user_ans,
            'correct_answer': q['answer'],
            'correct_answer_normalized': corr_norm,
            'correct': correct,
            'time': duration,
            'explanation': q.get('explanation','')
        })
        print(f"Temps: {duration:.2f}s")
    # Auto-évaluation Q open
    for r in quiz_results:
        if r['type']=='open' and r['correct'] is None:
            resp = input(f"Votre réponse était-elle correcte pour '{r['question']}' ? (o/n): ")
            r['correct'] = resp.lower().startswith('o')
    return quiz_results, total_time


def calculate_retention_metrics(quiz_results, total_time):
    """Métriques Ebbinghaus"""
    if not quiz_results: return {}
    valid = [r for r in quiz_results if r['correct'] is not None]
    correct_count = sum(r['correct'] for r in valid)
    total_q = len(valid)
    score = (correct_count/total_q*100) if total_q else 0
    avg_time = total_time/total_q if total_q else 0
    if score>=90: level='Excellent'
    elif score>=75: level='Bon'
    elif score>=50: level='Moyen'
    else: level='Faible'
    strength = score/100*max(0.5, min(1.0,15/max(1,avg_time)))
    target=0.7
    try:
        hours = -strength*100*math.log(target)
    except:
        hours = 24
    now = datetime.now()
    if hours<1:
        rev = f"dans {round(hours*60)}min"
    elif hours<24:
        rev = f"dans {round(hours)}h"
    else:
        days=round(hours/24)
        rev = f"dans {days}j"
    # Espacement
    reps=[hours, hours*2, hours*5, hours*12]
    rep_fmt=[]
    for h in reps:
        if h<24:
            rep_fmt.append(f"{round(h)}h")
        else:
            rep_fmt.append(f"{round(h/24)}j")
    return {
        'score': round(score,1),
        'correct_count': correct_count,
        'total_questions': total_q,
        'avg_time': avg_time,
        'confidence_level': level,
        'next_review': rev,
        'memory_strength': round(strength*100,2),
        'forgetting_rate': round((1-strength)*100,2),
        'spaced_repetition': rep_fmt
    }


def display_quiz_results(results, metrics):
    """Affiche résultats et métriques"""
    print("===== RÉSULTATS =====")
    print(f"Score: {metrics['score']}% ({metrics['correct_count']}/{metrics['total_questions']})")
    print(f"Temps moyen: {metrics['avg_time']:.2f}s")
    print(f"Confiance: {metrics['confidence_level']}")
    print(f"Prochaine révision: {metrics['next_review']}")
    print("Planning:")
    for i,rev in enumerate(metrics['spaced_repetition'],1): print(f" {i}. {rev}")


def save_results(content, base_name, formats, quiz_results=None, retention_metrics=None):
    """Sauvegarde TXT/PDF"""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    results = {}
    if 'txt' in formats:
        txt_file = f"{base_name}_{ts}.txt"
        try:
            with open(txt_file,'w',encoding='utf-8') as f:
                f.write(content)
                if quiz_results and retention_metrics:
                    f.write(f"\nScore: {retention_metrics['score']}%\n")
            results['txt'] = txt_file
        except Exception as e:
            print(f"Erreur TXT: {e}")
    if 'pdf' in formats:
        pdf_file = f"{base_name}_{ts}.pdf"
        try:
            pdf = SmartPDF()
            pdf.add_page()
            section = content.split('\n',1)[-1]
            pdf.add_section("Contenu", section)
            if quiz_results and retention_metrics:
                pdf.add_page()
                pdf.add_section("Résultats Quiz", str(retention_metrics))
            pdf.output(pdf_file)
            results['pdf'] = pdf_file
        except Exception as e:
            print(f"Erreur PDF: {e}")
    return results


def main():
    print("=== MemoryMesh - Analyse ===")
    print("1. Résumé\n2. Quiz\n3. Les deux")
    choice = input("Choix (1-3):").strip()
    if choice not in ('1','2','3'): return
    path = input("Fichier (PDF/PPTX/TXT):").strip('"')
    if not os.path.exists(path): return
    text = extract_text(path)
    if not text: return
    base = os.path.splitext(os.path.basename(path))[0]
    results, quiz_res, metrics = {}, None, None
    if choice in ('1','3'):
        summ = analyze_content(text,'summary')
        if summ:
            print("\n=== RESUMÉ ===\n", summ)
            results.update(save_results(f"Résumé\n{summ}", f"resume_{base}", ['txt','pdf']))
    if choice in ('2','3'):
        quiz_txt = analyze_content(text,'quiz')
        if quiz_txt:
            print("\nQuiz généré.")
            quiz_res, _ = parse_quiz(quiz_txt), None
            quiz_res, ttime = present_quiz(quiz_res)
            metrics = calculate_retention_metrics(quiz_res, ttime)
            display_quiz_results(quiz_res, metrics)
            results.update(save_results(quiz_txt, f"quiz_{base}", ['txt','pdf'], quiz_res, metrics))
    if results:
        print("Résultats:")
        for fmt,path in results.items(): print(f" - {fmt}: {os.path.abspath(path)}")

if __name__ == "__main__":
    required = {'requests','PyPDF2','python-pptx','fpdf'}
    installed = {pkg.split('==')[0] for pkg in os.popen('pip freeze').read().splitlines()}
    if missing := required - installed:
        os.system(f"pip install {' '.join(missing)}")
    main()
